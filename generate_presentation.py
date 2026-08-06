#!/usr/bin/env python3
"""
Generate ResumePilot AI Presentation (PPTX + PDF)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ──────────────────────────────────────────
PRIMARY   = RGBColor(0x1E, 0x40, 0x6E)  # Dark navy
ACCENT    = RGBColor(0x2E, 0x86, 0xDE)  # Bright blue
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)  # Teal/green
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)  # Very dark
BG_CARD   = RGBColor(0x1A, 0x23, 0x3B)  # Card dark
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
MUTED     = RGBColor(0x88, 0x99, 0xAA)
GOLD      = RGBColor(0xF5, 0xA6, 0x23)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ── Helper Functions ─────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    return add_shape(slide, left, top, width, height, fill_color=color)


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=LIGHT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_icon_circle(slide, left, top, size, color=ACCENT, label="", label_size=12):
    """Add a circle with a label beneath it."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if label:
        add_text_box(slide, left - Inches(0.3), top + size + Pt(4), size + Inches(0.6), Inches(0.5),
                     label, font_size=label_size, color=LIGHT, alignment=PP_ALIGN.CENTER)
    return shape


def add_numbered_step(slide, left, top, number, title, desc, color=ACCENT):
    """A numbered step with circle + title + description."""
    circle_size = Inches(0.6)
    add_icon_circle(slide, left, top, circle_size, color=color,
                    label="", label_size=10)
    # Number inside circle
    add_text_box(slide, left, top + Pt(2), circle_size, circle_size,
                 str(number), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    txBox = add_text_box(slide, left + circle_size + Inches(0.2), top - Pt(2),
                         Inches(4.5), Inches(0.4), title,
                         font_size=16, color=WHITE, bold=True)
    add_text_box(slide, left + circle_size + Inches(0.2), top + Inches(0.35),
                 Inches(4.5), Inches(0.8), desc,
                 font_size=12, color=MUTED)


# ── SLIDE 1: TITLE ───────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent bar top
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "ResumePilot AI", font_size=54, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "AI-Powered Resume Optimisation & ATS Scoring Engine",
             font_size=26, color=ACCENT, bold=False)

# Tagline
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Stop guessing. Start optimising.",
             font_size=18, color=MUTED)

# Bottom bar
add_accent_bar(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), ACCENT2)

# Presenter
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Abdu Ali Adem  |  Kanz AI Hackathon  |  2026",
             font_size=14, color=MUTED)


# ── SLIDE 2: THE PROBLEM ─────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "The Problem", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
             "Job seekers are flying blind into the biggest hiring filter they face.",
             font_size=20, color=LIGHT)

items = [
    "🔒 90% of large companies use Applicant Tracking Systems (ATS) to filter resumes",
    "❌ Qualified candidates get rejected — not for lack of fit, but for missing keywords or formatting",
    "📉 No feedback loop — candidates submit dozens of applications and never know why they're rejected",
    "🕳️ The ATS is a black box: candidates have zero visibility into what the system is looking for",
    "⏳ Hours spent tailoring each resume with no data to guide the effort",
]
add_bullet_slide(slide, Inches(1), Inches(2.6), Inches(11), Inches(4), items, font_size=16, color=LIGHT)


# ── SLIDE 3: THE SOLUTION ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "The Solution", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT2)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
             "ResumePilot AI gives candidates the same visibility into hiring that recruiters already have.",
             font_size=20, color=LIGHT)

# Three pillars
pillar_data = [
    ("🔍", "Analyse", "Upload your resume and\na job description for\ninstant ATS analysis"),
    ("⚡", "Score", "Get a precise match\npercentage with keyword\ndensity breakdown"),
    ("🛠️", "Optimise", "Rewrite your resume,\ngenerate a cover letter,\nand export a PDF report"),
]
for i, (icon, title, desc) in enumerate(pillar_data):
    x = Inches(1.2 + i * 4.0)
    y = Inches(2.8)
    # Card background
    add_shape(slide, x, y, Inches(3.4), Inches(3.4), fill_color=BG_CARD, line_color=RGBColor(0x2A, 0x35, 0x50), line_width=Pt(1))
    # Icon
    add_text_box(slide, x + Inches(1.2), y + Inches(0.3), Inches(1), Inches(0.8),
                 icon, font_size=40, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(3), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(3), Inches(1.2),
                 desc, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── SLIDE 4: HOW IT WORKS ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "How It Works", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

steps = [
    ("1", "Upload & Parse",
     "Upload your resume (PDF, DOCX, or TXT). The app extracts and cleans the text using pdfplumber or python-docx."),
    ("2", "AI Analysis",
     "Gemini 2.5 Flash performs semantic analysis — understanding intent, not just keyword counting. Combined with a custom scoring engine."),
    ("3", "Results Dashboard",
     "View your ATS match %, keyword density chart, missing skills, interview tips, and detailed scoring breakdown."),
    ("4", "Actionable Outputs",
     "Get a rewritten resume with missing keywords injected, a tailored cover letter, career coaching tips, and a downloadable PDF summary."),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.8 + i * 1.3)
    add_numbered_step(slide, Inches(1.2), y, num, title, desc)


# ── SLIDE 5: KEY FEATURES ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "Key Features", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT2)

features = [
    ("ATS Match Score", "Precise percentage with\nsemantic understanding"),
    ("Keyword Density Chart", "Visual breakdown of\nkeyword overlap vs. gaps"),
    ("Missing Skills", "AI-detected skills you're\nmissing for the role"),
    ("Interview Tips", "Tailored questions based\non job requirements"),
    ("Resume Rewrite", "Missing keywords injected\nnaturally into your CV"),
    ("Cover Letter", "AI-generated cover letter\nspecific to the job"),
    ("PDF Export", "Downloadable executive\nsummary report"),
    ("Job Comparison", "Side-by-side comparison\nof multiple job descriptions"),
]
for i, (title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(1.8 + row * 2.5)
    add_shape(slide, x, y, Inches(2.7), Inches(2.0), fill_color=BG_CARD, line_color=RGBColor(0x2A, 0x35, 0x50), line_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.3), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.3), Inches(1.0),
                 desc, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── SLIDE 6: TECH STACK ──────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "Technology Stack", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

tech_categories = [
    ("Frontend / UI", ["Streamlit", "Plotly Express", "HTML/CSS"], ACCENT),
    ("AI Engine", ["Google Gemini 2.5 Flash", "Prompt Engineering", "Semantic Analysis"], ACCENT2),
    ("Backend / Parsing", ["Python", "pdfplumber", "python-docx", "Pandas"], GOLD),
    ("Output / Export", ["ReportLab (PDF)", "Natural Language Generation"], ACCENT),
    ("DevOps / Tools", ["Git", "GitHub", "Streamlit Cloud", "VS Code"], LIGHT),
]

for i, (cat, items, color) in enumerate(tech_categories):
    x = Inches(0.8 + i * 2.5)
    y = Inches(1.8)
    # Card
    add_shape(slide, x, y, Inches(2.2), Inches(4.8), fill_color=BG_CARD, line_color=RGBColor(0x2A, 0x35, 0x50), line_width=Pt(1))
    # Category header
    add_shape(slide, x, y, Inches(2.2), Inches(0.6), fill_color=color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.0), Inches(0.5),
                 cat, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Items
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8 + j * 0.5), Inches(1.8), Inches(0.4),
                     f"• {item}", font_size=12, color=LIGHT)


# ── SLIDE 7: RESULTS & IMPACT ────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "Results & Impact", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT2)

# Impact metrics
metrics = [
    ("90%", "ATS pass rate\nimprovement target"),
    ("9", "Actionable views\nper analysis"),
    ("< 10s", "Analysis time\nper resume"),
    ("100%", "Free to use\nfor job seekers"),
]
for i, (num, label) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    add_shape(slide, x, y, Inches(2.7), Inches(2.2), fill_color=BG_CARD, line_color=RGBColor(0x2A, 0x35, 0x50), line_width=Pt(1))
    add_text_box(slide, x, y + Inches(0.3), Inches(2.7), Inches(0.8),
                 num, font_size=40, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.2), Inches(2.7), Inches(0.8),
                 label, font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
             "Built for job seekers, career coaches, and university career centres.\n"
             "Transforms the application process from a guessing game into a data-driven exercise.",
             font_size=15, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ── SLIDE 8: FUTURE ROADMAP ──────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "Future Roadmap", font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

roadmap = [
    ("📊", "Resume Version History", "Track score improvements over multiple iterations"),
    ("🔍", "Job Recommendation Engine", "Analyse your resume and suggest matching roles from live job boards"),
    ("👤", "LinkedIn Profile Analyser", "Extend optimisation beyond resumes to full professional presence"),
    ("🔄", "SDK Migration", "Upgrade from deprecated google.generativeai to the new GenAI SDK"),
    ("🎓", "Academic CV Mode", "Tailored version for academic CVs and grant applications"),
    ("💼", "Recruitment Side Tool", "Help hiring teams write better job descriptions based on top-scoring resumes"),
    ("🔐", "Subscription Tier", "Supabase auth + Stripe billing — already prototyped"),
]
for i, (icon, title, desc) in enumerate(roadmap):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 1.7)
    add_shape(slide, x, y, Inches(3.6), Inches(1.4), fill_color=BG_CARD, line_color=RGBColor(0x2A, 0x35, 0x50), line_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=WHITE)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.1), Inches(2.7), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.55), Inches(2.7), Inches(0.7),
                 desc, font_size=11, color=MUTED)


# ── SLIDE 9: CLOSING ─────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
             "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             "Try ResumePilot AI today", font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
             "https://resumepilot-ai-vngmvb9m6rdgszr7bthtbk.streamlit.app/",
             font_size=16, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
             "Abdu Ali Adem  |  abduaali132012@gmail.com",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Kanz AI Hackathon  |  2026",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_accent_bar(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), ACCENT2)


# ── SAVE PPTX ────────────────────────────────────────────────
pptx_path = "/app/ResumePilot_AI_Presentation.pptx"
prs.save(pptx_path)
print(f"✅ PPTX saved: {pptx_path}")


# ── GENERATE PDF ─────────────────────────────────────────────
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.units import inch, mm
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT

pdf_path = "/app/ResumePilot_AI_Presentation.pdf"

# Colours
C_PRIMARY   = HexColor("#1E406E")
C_ACCENT    = HexColor("#2E86DE")
C_ACCENT2   = HexColor("#10B981")
C_BG_DARK   = HexColor("#0F172A")
C_BG_CARD   = HexColor("#1A233B")
C_WHITE     = HexColor("#FFFFFF")
C_LIGHT     = HexColor("#CCDDEE")
C_MUTED     = HexColor("#8899AA")
C_GOLD      = HexColor("#F5A623")

doc = SimpleDocTemplate(
    pdf_path,
    pagesize=landscape(A4),
    topMargin=0.5*inch,
    bottomMargin=0.5*inch,
    leftMargin=0.5*inch,
    rightMargin=0.5*inch,
    title="ResumePilot AI - Presentation",
    author="Abdu Ali Adem",
)

styles = getSampleStyleSheet()

# Custom styles
s_title = ParagraphStyle("Title2", parent=styles["Title"], fontSize=36, textColor=C_WHITE, spaceAfter=12, alignment=TA_CENTER, fontName="Helvetica-Bold")
s_subtitle = ParagraphStyle("Sub2", parent=styles["Normal"], fontSize=20, textColor=C_ACCENT, spaceAfter=6, alignment=TA_CENTER, fontName="Helvetica")
s_body = ParagraphStyle("Body2", parent=styles["Normal"], fontSize=13, textColor=C_LIGHT, spaceAfter=8, leading=18, fontName="Helvetica")
s_body_bold = ParagraphStyle("BodyBold", parent=s_body, textColor=C_WHITE, fontName="Helvetica-Bold")
s_bullet = ParagraphStyle("Bullet", parent=s_body, leftIndent=20, bulletIndent=0, spaceBefore=4, spaceAfter=4)
s_section = ParagraphStyle("Section", parent=styles["Normal"], fontSize=28, textColor=C_WHITE, spaceAfter=10, spaceBefore=10, fontName="Helvetica-Bold")
s_metric_num = ParagraphStyle("MetricNum", parent=styles["Normal"], fontSize=36, textColor=C_ACCENT, alignment=TA_CENTER, fontName="Helvetica-Bold")
s_metric_label = ParagraphStyle("MetricLabel", parent=styles["Normal"], fontSize=12, textColor=C_MUTED, alignment=TA_CENTER, fontName="Helvetica")
s_card_title = ParagraphStyle("CardTitle", parent=styles["Normal"], fontSize=14, textColor=C_WHITE, alignment=TA_CENTER, fontName="Helvetica-Bold")
s_card_desc = ParagraphStyle("CardDesc", parent=styles["Normal"], fontSize=11, textColor=C_MUTED, alignment=TA_CENTER, fontName="Helvetica")
s_footer = ParagraphStyle("Footer", parent=styles["Normal"], fontSize=10, textColor=C_MUTED, alignment=TA_CENTER, fontName="Helvetica")

elements = []

# ── Helper to add a background rectangle ──
def add_bg_rect(canvas, doc):
    """Draw dark background on every page."""
    canvas.saveState()
    canvas.setFillColor(C_BG_DARK)
    w, h = landscape(A4)
    canvas.rect(0, 0, w, h, fill=1, stroke=0)
    # Top accent bar
    canvas.setFillColor(C_ACCENT)
    canvas.rect(0, h - 6, w, 6, fill=1, stroke=0)
    canvas.restoreState()

# ── PAGE 1: Title ──
elements.append(Spacer(1, 1.5*inch))
elements.append(Paragraph("ResumePilot AI", s_title))
elements.append(Spacer(1, 0.2*inch))
elements.append(Paragraph("AI-Powered Resume Optimisation<br/>&amp; ATS Scoring Engine", s_subtitle))
elements.append(Spacer(1, 0.4*inch))
elements.append(Paragraph("Stop guessing. Start optimising.", ParagraphStyle("Tagline", parent=s_body, fontSize=16, textColor=C_MUTED, alignment=TA_CENTER)))
elements.append(Spacer(1, 1.5*inch))
elements.append(Paragraph("Abdu Ali Adem  |  Kanz AI Hackathon  |  2026", s_footer))
elements.append(PageBreak())

# ── PAGE 2: The Problem ──
elements.append(Paragraph("The Problem", s_section))
elements.append(Spacer(1, 0.1*inch))
elements.append(Paragraph("Job seekers are flying blind into the biggest hiring filter they face.", ParagraphStyle("Lead", parent=s_body, fontSize=16, textColor=C_LIGHT, spaceAfter=12)))
bullets = [
    "&bull; 90% of large companies use Applicant Tracking Systems (ATS) to filter resumes",
    "&bull; Qualified candidates get rejected — not for lack of fit, but for missing keywords",
    "&bull; No feedback loop — candidates never know why they're rejected",
    "&bull; The ATS is a black box with zero visibility into what it's looking for",
    "&bull; Hours spent tailoring each resume with no data to guide the effort",
]
for b in bullets:
    elements.append(Paragraph(b, s_bullet))
elements.append(PageBreak())

# ── PAGE 3: The Solution ──
elements.append(Paragraph("The Solution", s_section))
elements.append(Paragraph("ResumePilot AI gives candidates the same visibility into hiring that recruiters already have.", ParagraphStyle("Lead2", parent=s_body, fontSize=16, textColor=C_LIGHT, spaceAfter=16)))

pillar_data = [
    ("🔍", "Analyse", "Upload your resume and a job description for instant ATS analysis"),
    ("⚡", "Score", "Get a precise match percentage with keyword density breakdown"),
    ("🛠️", "Optimise", "Rewrite your resume, generate a cover letter, and export a PDF report"),
]
pillar_rows = []
for icon, title, desc in pillar_data:
    cell = f'<para><b><font size="18" color="{C_WHITE.hexval()}">{icon}</font></b><br/><font size="14" color="{C_WHITE.hexval()}"><b>{title}</b></font><br/><font size="11" color="{C_MUTED.hexval()}">{desc}</font></para>'
    pillar_rows.append([cell])

t = Table(pillar_rows, colWidths=[3.8*inch]*3, rowHeights=[2.2*inch]*3)
t.setStyle(TableStyle([
    ('BACKGROUND', (0,0), (-1,-1), C_BG_CARD),
    ('BOX', (0,0), (-1,-1), 0.5, HexColor("#2A3350")),
    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
    ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ('TOPPADDING', (0,0), (-1,-1), 12),
    ('BOTTOMPADDING', (0,0), (-1,-1), 12),
    ('LEFTPADDING', (0,0), (-1,-1), 8),
    ('RIGHTPADDING', (0,0), (-1,-1), 8),
]))
elements.append(t)
elements.append(PageBreak())

# ── PAGE 4: How It Works ──
elements.append(Paragraph("How It Works", s_section))
elements.append(Spacer(1, 0.1*inch))
steps = [
    ("<b>1. Upload &amp; Parse</b>", "Upload your resume (PDF, DOCX, or TXT). The app extracts and cleans the text using pdfplumber or python-docx."),
    ("<b>2. AI Analysis</b>", "Gemini 2.5 Flash performs semantic analysis — understanding intent, not just keyword counting. Combined with a custom scoring engine."),
    ("<b>3. Results Dashboard</b>", "View your ATS match %, keyword density chart, missing skills, interview tips, and detailed scoring breakdown."),
    ("<b>4. Actionable Outputs</b>", "Get a rewritten resume with missing keywords injected, a tailored cover letter, career coaching tips, and a downloadable PDF summary."),
]
for title, desc in steps:
    elements.append(Paragraph(title, s_body_bold))
    elements.append(Paragraph(desc, s_bullet))
    elements.append(Spacer(1, 0.08*inch))
elements.append(PageBreak())

# ── PAGE 5: Key Features ──
elements.append(Paragraph("Key Features", s_section))
elements.append(Spacer(1, 0.1*inch))

features = [
    ("ATS Match Score", "Precise percentage with semantic understanding"),
    ("Keyword Density Chart", "Visual breakdown of keyword overlap vs. gaps"),
    ("Missing Skills", "AI-detected skills you're missing for the role"),
    ("Interview Tips", "Tailored questions based on job requirements"),
    ("Resume Rewrite", "Missing keywords injected naturally into your CV"),
    ("Cover Letter", "AI-generated cover letter specific to the job"),
    ("PDF Export", "Downloadable executive summary report"),
    ("Job Comparison", "Side-by-side comparison of multiple job descriptions"),
]

feat_rows = []
for i in range(0, len(features), 4):
    row = []
    for j in range(4):
        if i + j < len(features):
            title, desc = features[i + j]
            cell = f'<para><font size="13" color="{C_WHITE.hexval()}"><b>{title}</b></font><br/><font size="10" color="{C_MUTED.hexval()}">{desc}</font></para>'
            row.append(cell)
        else:
            row.append("")
    feat_rows.append(row)

t = Table(feat_rows, colWidths=[2.8*inch]*4, rowHeights=[1.2*inch]*2)
t.setStyle(TableStyle([
    ('BACKGROUND', (0,0), (-1,-1), C_BG_CARD),
    ('BOX', (0,0), (-1,-1), 0.5, HexColor("#2A3350")),
    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
    ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ('TOPPADDING', (0,0), (-1,-1), 10),
    ('BOTTOMPADDING', (0,0), (-1,-1), 10),
    ('LEFTPADDING', (0,0), (-1,-1), 6),
    ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ('GRID', (0,0), (-1,-1), 0.5, HexColor("#2A3350")),
]))
elements.append(t)
elements.append(PageBreak())

# ── PAGE 6: Tech Stack ──
elements.append(Paragraph("Technology Stack", s_section))
elements.append(Spacer(1, 0.1*inch))

tech_categories = [
    ("Frontend / UI", "Streamlit, Plotly Express, HTML/CSS"),
    ("AI Engine", "Google Gemini 2.5 Flash, Prompt Engineering, Semantic Analysis"),
    ("Backend / Parsing", "Python, pdfplumber, python-docx, Pandas"),
    ("Output / Export", "ReportLab (PDF), Natural Language Generation"),
    ("DevOps / Tools", "Git, GitHub, Streamlit Cloud, VS Code"),
]
for cat, items in tech_categories:
    elements.append(Paragraph(f'<font color="{C_ACCENT.hexval()}"><b>{cat}</b></font>', s_body))
    elements.append(Paragraph(f'&nbsp;&nbsp;&nbsp;{items}', s_bullet))
    elements.append(Spacer(1, 0.06*inch))
elements.append(PageBreak())

# ── PAGE 7: Future Roadmap ──
elements.append(Paragraph("Future Roadmap", s_section))
elements.append(Spacer(1, 0.1*inch))
roadmap = [
    ("📊 Resume Version History", "Track score improvements over multiple iterations"),
    ("🔍 Job Recommendation Engine", "Analyse your resume and suggest matching roles from live job boards"),
    ("👤 LinkedIn Profile Analyser", "Extend optimisation beyond resumes to full professional presence"),
    ("🔄 SDK Migration", "Upgrade from deprecated google.generativeai to the new GenAI SDK"),
    ("🎓 Academic CV Mode", "Tailored version for academic CVs and grant applications"),
    ("💼 Recruitment Side Tool", "Help hiring teams write better job descriptions"),
    ("🔐 Subscription Tier", "Supabase auth + Stripe billing — already prototyped"),
]
for title, desc in roadmap:
    elements.append(Paragraph(f'<font color="{C_WHITE.hexval()}"><b>{title}</b></font>', s_body))
    elements.append(Paragraph(f'&nbsp;&nbsp;&nbsp;{desc}', s_bullet))
    elements.append(Spacer(1, 0.04*inch))
elements.append(PageBreak())

# ── PAGE 8: Closing ──
elements.append(Spacer(1, 2*inch))
elements.append(Paragraph("Thank You", s_title))
elements.append(Spacer(1, 0.3*inch))
elements.append(Paragraph("Try ResumePilot AI today", ParagraphStyle("CTA", parent=s_subtitle, fontSize=18, textColor=C_ACCENT2, alignment=TA_CENTER)))
elements.append(Spacer(1, 0.2*inch))
elements.append(Paragraph("https://resumepilot-ai-vngmvb9m6rdgszr7bthtbk.streamlit.app/", ParagraphStyle("Link", parent=s_body, fontSize=14, textColor=C_ACCENT, alignment=TA_CENTER)))
elements.append(Spacer(1, 0.5*inch))
elements.append(Paragraph("Abdu Ali Adem  |  abduaali132012@gmail.com", s_footer))
elements.append(Paragraph("Kanz AI Hackathon  |  2026", s_footer))

# Build PDF
doc.build(elements, onFirstPage=add_bg_rect, onLaterPages=add_bg_rect)
print(f"✅ PDF saved: {pdf_path}")